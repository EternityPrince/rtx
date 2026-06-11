import pytest
import sys
import runpy
from pathlib import Path
from unittest.mock import patch, MagicMock, PropertyMock
from rich.console import Console

from rtx.digger import (
    DiggerEngine,
    is_binary_content,
    is_excluded_path,
    get_file_metrics,
    format_bytes
)
from rtx.parsers.base import BaseParser
from rtx.parsers.code import CodeParser
from rtx.parsers.text import TextParser
from rtx.parsers.docx import DocxParser, iter_block_items
from rtx.parsers.pptx import PptxParser
from rtx.parsers.pdf import PdfParser, get_marker_models, clear_marker_models
from rtx.parsers.ebook import EbookParser
from rtx.parsers import get_parser_for_path
from rtx.pipeline import run_parser_pipeline, ParseResult
from rtx.writer import write_stream, write_mirror, update_gitignore
from rtx.cli import app, print_summary_table

# ==========================================
# 1. BaseParser & CLI & runpy Coverage
# ==========================================

def test_base_parser():
    class DummyParser(BaseParser):
        def parse(self, path: Path) -> str:
            return super().parse(path)
            
    parser = DummyParser()
    assert parser.parse(Path(".")) is None

def test_summary_table_edge_cases():
    console = Console(record=True)
    
    # Empty results
    print_summary_table([], console)
    out_empty = console.export_text()
    assert "No files to process" in out_empty
    
    # Mixed and failing results
    results = [
        ParseResult(Path("no_ext"), "", "content", "Success"),
        ParseResult(Path("fail.py"), ".py", "", "Failed", "Fatal error"),
        ParseResult(Path("success.py"), ".py", "print(1)", "Success"),
        ParseResult(Path("fail.txt"), ".txt", "", "Failed", "Some error"),
    ]
    print_summary_table(results, console)
    out_mixed = console.export_text()
    assert "no_ext" in out_mixed
    assert "Partial" in out_mixed
    assert "Error" in out_mixed

def test_cli_tui_and_empty_directory_calls(sandbox_dir):
    from typer.testing import CliRunner
    
    # Test typer TUI call branch
    runner = CliRunner()
    with patch("rtx.tui.run_tui") as mock_run_tui:
        result = runner.invoke(app, ["--path", str(sandbox_dir)])
        assert result.exit_code == 0
        mock_run_tui.assert_called_once()
        
    # Test empty scanning directory (no valid files)
    empty_scan_dir = sandbox_dir / "empty_dir_test"
    empty_scan_dir.mkdir(exist_ok=True)
    (empty_scan_dir / "node_modules").mkdir(exist_ok=True)
    
    result2 = runner.invoke(app, ["scan", "--project", str(empty_scan_dir)])
    assert result2.exit_code == 0
    assert "No valid files to process" in result2.output

def test_cli_runpy_main():
    with patch.object(sys, "argv", ["rtx", "--help"]):
        with pytest.raises(SystemExit) as exc:
            runpy.run_module("rtx.cli", run_name="__main__")
        assert exc.value.code == 0

# ==========================================
# 2. Backtick Fence Collisions in Parsers
# ==========================================

def test_code_parser_fence_collision(sandbox_dir):
    file_path = sandbox_dir / "collision.py"
    file_path.write_text("print('hello')\n```\n# some block\n```\n", encoding="utf-8")
    
    parser = CodeParser()
    parsed = parser.parse(file_path)
    assert parsed.startswith("````python")
    assert parsed.endswith("````")
    file_path.unlink()

def test_text_parser_fence_collision(sandbox_dir):
    file_path = sandbox_dir / "collision.txt"
    file_path.write_text("Some text with ``` triple backticks", encoding="utf-8")
    
    parser = TextParser()
    parsed = parser.parse(file_path)
    assert parsed.startswith("````text")
    assert parsed.endswith("````")
    file_path.unlink()

# ==========================================
# 3. Docx Parser Lists, Tables & Exceptions
# ==========================================

def test_docx_parser_detailed_features(sandbox_dir):
    import docx
    doc = docx.Document()
    doc.add_heading("H1", level=1)
    doc.add_heading("H2", level=2)
    doc.add_heading("H3", level=3)
    doc.add_heading("H4", level=4)
    doc.add_paragraph("Bullet item", style="List Bullet")
    doc.add_paragraph("Numbered item", style="List Number")
    
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Col 1 Header"
    table.cell(0, 1).text = "Col 2 Header"
    table.cell(1, 0).text = "Row 1 Col 1"
    table.cell(1, 1).text = "Row 1 Col 2"
    
    docx_file = sandbox_dir / "complex.docx"
    doc.save(docx_file)
    
    parser = DocxParser()
    text = parser.parse(docx_file)
    
    assert "# H1" in text
    assert "## H2" in text
    assert "### H3" in text
    assert "#### H4" in text
    assert "* Bullet item" in text
    assert "1. Numbered item" in text
    assert "| Col 1 Header | Col 2 Header |" in text
    assert "| --- | --- |" in text
    assert "| Row 1 Col 1 | Row 1 Col 2 |" in text
    
    docx_file.unlink()

def test_iter_block_items_exceptions():
    from docx.table import Table
    mock_tbl_el = MagicMock()
    mock_parent = MagicMock()
    table_obj = Table(mock_tbl_el, mock_parent)
    table_obj._tbl = MagicMock()
    table_obj._tbl.iterchildren.return_value = []
    
    list(iter_block_items(table_obj))
    
    with pytest.raises(TypeError):
        list(iter_block_items("Unsupported String Parent"))

# ==========================================
# 4. Pptx Parser Features & Exceptions
# ==========================================

def test_pptx_parser_edge_cases(sandbox_dir):
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    slide.shapes.title.text = "Slide Title"
    
    tb = slide.shapes.add_textbox(100, 100, 100, 100)
    p_empty = tb.text_frame.add_paragraph()
    p_empty.text = ""
    
    from pptx.util import Inches
    slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(1), Inches(1))
    
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    tb2 = slide2.shapes.add_textbox(100, 100, 100, 100)
    tb2.text_frame.text = "Single Level 0 Text"
    
    pptx_file = sandbox_dir / "complex.pptx"
    prs.save(pptx_file)
    
    parser = PptxParser()
    text = parser.parse(pptx_file)
    
    assert "Slide Title" in text
    assert "Single Level 0 Text" in text
    pptx_file.unlink()
    
    with pytest.raises(Exception):
        parser.parse(Path("invalid_non_existent.pptx"))

# ==========================================
# 5. PDF Parser Models & Fallbacks
# ==========================================

def test_pdf_models_caching_and_cleanup():
    mock_models = MagicMock()
    mock_models.create_model_dict.return_value = "mocked_models"
    
    with patch.dict(sys.modules, {"marker.models": mock_models}):
        # Force cache reload by resetting module global
        with patch("rtx.parsers.pdf._marker_models", None):
            models = get_marker_models()
            assert models == "mocked_models"
            
            models2 = get_marker_models()
            assert models2 == "mocked_models"
            
            # Test clear_marker_models with torch loaded
            mock_torch = MagicMock()
            mock_torch.cuda.is_available.return_value = True
            mock_torch.mps.is_available.return_value = True
            
            with patch.dict("sys.modules", {"torch": mock_torch}):
                clear_marker_models()
                mock_torch.cuda.empty_cache.assert_called_once()
                mock_torch.mps.empty_cache.assert_called_once()
                
            # Test clear_marker_models with torch Exception
            mock_torch_err = MagicMock()
            mock_torch_err.cuda.is_available.side_effect = Exception("error")
            with patch.dict("sys.modules", {"torch": mock_torch_err}):
                clear_marker_models()  # should handle exception silently
                
    assert get_marker_models.__globals__.get("_marker_models") is None

def test_pdf_parser_marker_success(sandbox_dir):
    mock_models = MagicMock()
    mock_models.create_model_dict.return_value = "mock_models"
    mock_converter_inst = MagicMock()
    mock_rendered = MagicMock()
    mock_rendered.markdown = "Marker Parsed Text"
    mock_converter_inst.return_value = mock_rendered
    mock_converter_class = MagicMock(return_value=mock_converter_inst)
    
    with patch.dict(sys.modules, {
        "marker.models": mock_models,
        "marker.converters.pdf": MagicMock(PdfConverter=mock_converter_class)
    }):
        with patch("rtx.parsers.pdf._marker_models", None):
            parser = PdfParser()
            text = parser.parse(sandbox_dir / "sample.pdf")
            assert text == "Marker Parsed Text"

def test_pdf_pypdf_fallback_success(sandbox_dir):
    with patch("rtx.parsers.pdf.get_marker_models", side_effect=ImportError):
        with patch("fitz.open", side_effect=Exception("Fitz Open Error")):
            mock_reader = MagicMock()
            mock_page1 = MagicMock()
            mock_page1.extract_text.return_value = "Page 1 PyPDF Text"
            mock_page2 = MagicMock()
            mock_page2.extract_text.return_value = None
            mock_reader.pages = [mock_page1, mock_page2]
            
            with patch("pypdf.PdfReader", return_value=mock_reader):
                parser = PdfParser()
                text = parser.parse(sandbox_dir / "sample.pdf")
                assert text == "Page 1 PyPDF Text"

def test_pdf_pypdf_fallback_empty(sandbox_dir):
    with patch("rtx.parsers.pdf.get_marker_models", side_effect=ImportError):
        with patch("fitz.open", side_effect=Exception("Fitz Open Error")):
            mock_reader = MagicMock()
            mock_page = MagicMock()
            mock_page.extract_text.return_value = ""
            mock_reader.pages = [mock_page]
            
            with patch("pypdf.PdfReader", return_value=mock_reader):
                parser = PdfParser()
                assert parser.parse(sandbox_dir / "sample.pdf") == ""

# ==========================================
# 6. Ebooks: EPUB, FB2 & Pipeline
# ==========================================

def test_ebook_parser_unsupported():
    parser = EbookParser()
    with pytest.raises(ValueError):
        parser.parse(Path("file.mobi"))

def test_fb2_and_pipeline_full_batch(sandbox_dir):
    fb2_file = sandbox_dir / "book.fb2"
    fb2_file.write_text(
        "<fictionbook><title-info><book-title>BookTitle</book-title></title-info>"
        "<body><section><title>Chapter Header</title><p>Para</p></section></body></fictionbook>",
        encoding="utf-8"
    )
    
    # Run pipeline with FB2 included
    results = run_parser_pipeline([fb2_file])
    assert len(results) == 1
    assert results[0].status == "Success"
    assert "BookTitle" in results[0].text
    
    # Test get_file_metrics on FB2
    lines, chars, bytes_sz = get_file_metrics(fb2_file)
    assert lines > 0
    assert chars > 0
    assert bytes_sz == fb2_file.stat().st_size
    
    fb2_file.unlink()
    
    # Test failed FB2 metrics calculation
    invalid_fb2 = Path("non_existent_file.fb2")
    assert get_file_metrics(invalid_fb2) == (0, 0, 0)

# ==========================================
# 7. Digger Engine Recursion & Edge Cases
# ==========================================

def test_digger_recursion_and_cache(sandbox_dir):
    # Create sub-directory structure
    sub_dir = sandbox_dir / "sub_folder"
    sub_dir.mkdir(exist_ok=True)
    sub_file = sub_dir / "nested_code.py"
    sub_file.write_text("print('nested')", encoding="utf-8")
    
    digger = DiggerEngine(sandbox_dir)
    
    files = list(digger.scan_valid_files())
    assert sub_file in files
    
    total_files, total_lines, total_chars, total_bytes = digger.calculate_metrics(sandbox_dir)
    assert total_files > 0
    
    total_files2, _, _, _ = digger.calculate_metrics(sandbox_dir)
    assert total_files2 == total_files
    
    excl_metrics = digger.calculate_metrics(sandbox_dir / ".git")
    assert excl_metrics == (0, 0, 0, 0)
    
    file_with_read_error = sandbox_dir / "error_file.txt"
    file_with_read_error.write_text("test")
    with patch("builtins.open", side_effect=OSError("Read error")):
        lines_err, chars_err, bytes_err = get_file_metrics(file_with_read_error)
        assert lines_err == 0
        assert chars_err == 0
    file_with_read_error.unlink()

# ==========================================
# 8. Writer Edge Cases
# ==========================================

def test_writer_gitignore_and_copy_exceptions(sandbox_dir):
    gitignore = sandbox_dir / ".gitignore"
    
    # Case A: .gitignore exists, does not end with \n
    gitignore.write_text("# initial comment\n.env", encoding="utf-8")
    update_gitignore(sandbox_dir)
    assert gitignore.read_text(encoding="utf-8") == "# initial comment\n.env\n.rtx/\n"
    
    # Case B: .gitignore exists, ends with \n
    gitignore.write_text("# comment\n", encoding="utf-8")
    update_gitignore(sandbox_dir)
    assert gitignore.read_text(encoding="utf-8") == "# comment\n.rtx/\n"
    
    gitignore.unlink()
    
    with patch("builtins.open", side_effect=Exception("Permission Error")):
        update_gitignore(sandbox_dir)
        
    res = [ParseResult(sandbox_dir / "calculator.py", ".py", "def add(): pass", "Success")]
    with patch("pyperclip.copy", side_effect=Exception("Clipboard Unavailable")):
        write_stream(res, sandbox_dir, output_file=None)
        
    failed_res = [ParseResult(sandbox_dir / "failed.py", ".py", "", "Failed", "Some Error")]
    write_mirror(failed_res, sandbox_dir)
    assert not (sandbox_dir / ".rtx" / "failed.py.md").exists()
    
    # Case F: write_mirror path not relative to project root (ValueError fallback)
    outside_res = [ParseResult(Path("/outside/root/file.py"), ".py", "print(99)", "Success")]
    write_mirror(outside_res, sandbox_dir)
    assert (sandbox_dir / ".rtx" / "file.py.md").exists()
    
    with patch("builtins.open", side_effect=OSError("Write failure")):
        write_mirror(res, sandbox_dir)

# ==========================================
# 9. TUI Widget rendering & event actions
# ==========================================

@pytest.mark.asyncio
async def test_tui_interactive_actions(sandbox_dir):
    from rtx.tui import RtxApp, ParsingScreen
    from textual.widgets import Input, Button
    from unittest.mock import patch
    
    with patch("rtx.parsers.pdf.PdfParser.parse", return_value="Mocked PDF Content"):
        app = RtxApp(sandbox_dir)
        
        async with app.run_test() as pilot:
            # Toggle mirror mode using keyboard
            await pilot.press("m")
            assert app.mirror_mode is True
            await pilot.press("m")
            assert app.mirror_mode is False
            
            # Test input path updates
            await pilot.click("#output_path_input")
            app.query_one("#output_path_input", Input).value = "tui_results.md"
            await pilot.press("enter")
            assert app.output_file == Path("tui_results.md")
            
            # Clear output path
            app.query_one("#output_path_input", Input).value = ""
            await pilot.press("enter")
            assert app.output_file is None
            
            # Select files using toggle action on current node (root)
            app.action_toggle_select()
            assert len(app.selected_paths) > 0
            
            # Trigger parsing modal screen
            app.action_start_parsing()
            await pilot.pause()
            
            # Assert modal is mounted
            assert isinstance(app.screen, ParsingScreen)
            
            # Wait for background thread parsing to finish
            await pilot.pause(0.5)
            
            # Close button should be enabled/visible now
            close_btn = app.screen.query_one("#close_btn", Button)
            assert "hidden" not in close_btn.classes
            
            # Click close to return to tree
            await pilot.click("#close_btn")
            await pilot.pause()
            assert not isinstance(app.screen, ParsingScreen)
            
            # Quit TUI
            await pilot.press("q")

@pytest.mark.asyncio
async def test_tui_widget_coverage(sandbox_dir):
    from rtx.tui import RtxApp, ParsingScreen, RtxDirectoryTree, run_tui
    
    app = RtxApp(sandbox_dir)
    
    async with app.run_test() as pilot:
        # Prepopulate node stats to cover stats rendering
        app.node_stats[sandbox_dir] = (5, 120, 2400, 95000)
        
        tree = app.query_one(RtxDirectoryTree)
        tree.refresh()
        await pilot.pause()
        
        # Test render_label file stat Exception branch
        with patch("pathlib.Path.stat", side_effect=Exception("Stat failure")):
            tree.refresh()
            await pilot.pause()
            
        # Test checkbox statuses: selected and partially selected folders
        app.selected_paths = {sandbox_dir / "calculator.py"}
        tree.refresh()
        await pilot.pause()
        
        app.selected_paths = {sandbox_dir / "sub_folder" / "nested_code.py"}
        tree.refresh()
        await pilot.pause()
        
        # Test run_tui call
        with patch("rtx.tui.RtxApp.run") as mock_run:
            run_tui(sandbox_dir)
            mock_run.assert_called_once()
            
        # Test action_toggle_select when cursor_node is None
        with patch("textual.widgets.DirectoryTree.cursor_node", new_callable=PropertyMock) as mock_cursor:
            mock_cursor.return_value = None
            app.action_toggle_select()
            
        # Mock cursor node representing a file
        mock_node_file = MagicMock()
        mock_node_file.data.path = sandbox_dir / "calculator.py"
        
        with patch("textual.widgets.DirectoryTree.cursor_node", new_callable=PropertyMock) as mock_cursor:
            mock_cursor.return_value = mock_node_file
            
            # Select file
            app.selected_paths.clear()
            app.action_toggle_select()
            assert sandbox_dir / "calculator.py" in app.selected_paths
            
            # Unselect file
            app.action_toggle_select()
            assert sandbox_dir / "calculator.py" not in app.selected_paths
            
        # Mock cursor node representing a directory
        mock_node_dir = MagicMock()
        mock_node_dir.data.path = sandbox_dir
        
        with patch("textual.widgets.DirectoryTree.cursor_node", new_callable=PropertyMock) as mock_cursor:
            mock_cursor.return_value = mock_node_dir
            
            # Select all in directory
            app.selected_paths.clear()
            app.action_toggle_select()
            assert sandbox_dir in app.selected_paths
            
            # Unselect all in directory
            app.action_toggle_select()
            assert sandbox_dir not in app.selected_paths
            
        # Test starting parsing with empty selection (raises warning notification)
        app.selected_paths = set()
        app.action_start_parsing()
        await pilot.pause()
        
        # Test finished_parsing screen statuses
        pscreen = ParsingScreen(
            selected_files=[sandbox_dir / "calculator.py"],
            mirror_mode=True,
            output_file=None,
            project_path=sandbox_dir
        )
        
        app.push_screen(pscreen)
        await pilot.pause()
        
        mixed_results = [
            ParseResult(sandbox_dir / "ok.py", ".py", "ok", "Success"),
            ParseResult(sandbox_dir / "err.py", ".py", "", "Failed", "Exception detail")
        ]
        pscreen.finished_parsing(mixed_results)
        await pilot.pause()
        
        await pilot.press("q")

# ==========================================
# 10. Additional Gaps Coverage
# ==========================================

def test_digger_extra_coverage(sandbox_dir):
    # 1. is_binary_content raising OSError (file not found is OSError)
    assert is_binary_content(Path("non_existent_file.xyz")) is True

    # 2. is_excluded_path ValueError (outside root)
    assert is_excluded_path(Path("/outside/root/file.py"), sandbox_dir) is True

    # 3. is_excluded_path Exception (path is None)
    assert is_excluded_path(None, sandbox_dir) is True

    # 4. calculate_metrics with excluded path
    digger = DiggerEngine(sandbox_dir)
    assert digger.calculate_metrics(sandbox_dir / "image.png") == (0, 0, 0, 0)

    # 5. calculate_metrics scandir OSError
    with patch("os.scandir", side_effect=OSError("Access denied")):
        # Calculate metrics for sandbox_dir, should handle OSError and return (0,0,0,0) or partial
        digger.metrics_cache.clear()
        assert digger.calculate_metrics(sandbox_dir) == (0, 0, 0, 0)

    # 6. count_pdf_metrics fallback to pypdf (mock fitz to raise Exception)
    from rtx.digger import count_pdf_metrics
    with patch("fitz.open", side_effect=Exception("Fitz error")):
        lines, chars = count_pdf_metrics(sandbox_dir / "sample.pdf")
        assert lines > 0 or chars >= 0

    # 7. count_docx_metrics table cell paragraphs and exceptions
    from rtx.digger import count_docx_metrics, count_pptx_metrics, count_epub_metrics, count_fb2_metrics
    assert count_docx_metrics(Path("non_existent.docx")) == (0, 0)
    assert count_pptx_metrics(Path("non_existent.pptx")) == (0, 0)
    assert count_epub_metrics(Path("non_existent.epub")) == (0, 0)
    assert count_fb2_metrics(Path("non_existent.fb2")) == (0, 0)

    # DOCX table metrics test
    import docx
    doc = docx.Document()
    table = doc.add_table(rows=1, cols=1)
    table.cell(0, 0).text = "Table Cell Content"
    docx_file = sandbox_dir / "table_test.docx"
    doc.save(docx_file)
    lines, chars = count_docx_metrics(docx_file)
    assert chars > 0
    docx_file.unlink()


def test_pdf_parser_failure_branch(sandbox_dir):
    # Mock everything to fail in PdfParser.parse to raise RuntimeError
    with patch("rtx.parsers.pdf.get_marker_models", side_effect=Exception("marker error")):
        with patch("fitz.open", side_effect=Exception("fitz error")):
            with patch("pypdf.PdfReader", side_effect=Exception("pypdf error")):
                parser = PdfParser()
                with pytest.raises(RuntimeError) as exc:
                    parser.parse(sandbox_dir / "sample.pdf")
                assert "PDF parsing failed" in str(exc.value)


@pytest.mark.asyncio
async def test_tui_extra_coverage(sandbox_dir):
    from rtx.tui import RtxApp, ParsingScreen
    from textual.widgets import Button
    
    # Create a non-excluded subdirectory and a file in it
    sub_dir = sandbox_dir / "tui_sub_folder"
    sub_dir.mkdir(exist_ok=True)
    sub_file = sub_dir / "nested.py"
    sub_file.write_text("print('test')", encoding="utf-8")
    
    app = RtxApp(sandbox_dir)
    async with app.run_test() as pilot:
        # 1. Test action_reload
        app.node_stats[sandbox_dir] = (1, 1, 1, 1)
        app.digger.metrics_cache[sandbox_dir] = (1, 1, 1, 1)
        app.action_reload()
        assert not app.node_stats
        assert not app.digger.metrics_cache

        # 2. Test update_details with >15 files
        app.selected_paths = {sandbox_dir / f"file_{i}.py" for i in range(20)}
        app.update_details()

        # 3. Test partially selected folder prefix (🟨) in render_label
        app.selected_paths = {sub_file.resolve()}
        tree = app.query_one("#dir_tree")
        tree.refresh()
        await pilot.pause()

        # 4. Test action_open_file when current_preview_path is None or a directory
        app.current_preview_path = None
        app.action_open_file()
        
        app.current_preview_path = sandbox_dir
        app.action_open_file()

        # 5. Test open_current_file for different OS platforms
        app.current_preview_path = sub_file.resolve()
        with patch("subprocess.run") as mock_sub:
            with patch("sys.platform", "darwin"):
                app.open_current_file()
                mock_sub.assert_called_with(["open", str(sub_file.resolve())], check=True)
                
            with patch("sys.platform", "linux"):
                app.open_current_file()
                mock_sub.assert_called_with(["xdg-open", str(sub_file.resolve())], check=True)

        with patch("os.startfile", create=True) as mock_start:
            with patch("sys.platform", "win32"):
                app.open_current_file()
                mock_start.assert_called_with(sub_file.resolve())

        # Test exception inside open_current_file
        with patch("subprocess.run", side_effect=Exception("Subprocess error")):
            with patch("sys.platform", "darwin"):
                app.open_current_file()

        # 6. Test show_preview with large file (>10MB)
        large_file = sandbox_dir / "large.py"
        import stat
        mock_stat_obj = MagicMock()
        mock_stat_obj.st_size = 11 * 1024 * 1024
        mock_stat_obj.st_mode = stat.S_IFREG
        
        with patch("pathlib.Path.stat", return_value=mock_stat_obj):
            app.show_preview(large_file)
            await pilot.pause()
            
        # Trigger exception in path.stat()
        with patch("pathlib.Path.is_dir", return_value=False), \
             patch("pathlib.Path.is_file", return_value=True), \
             patch("pathlib.Path.stat", side_effect=Exception("stat error")):
            app.show_preview(large_file)
            await pilot.pause()

        # 7. Test show_preview with binary file content
        bin_file = sandbox_dir / "binary_test.py"
        bin_file.write_bytes(b"some text \x00 more text")
        app.show_preview(bin_file)
        await pilot.pause()
        bin_file.unlink()

        # 8. Test show_preview fallback when bat fails or is missing
        with patch("subprocess.run", side_effect=FileNotFoundError):
            app.show_preview(sub_file)
            await pilot.pause()

        # 9. Test FinishedParsing screen with all failed results
        pscreen = ParsingScreen(
            selected_files=[sub_file],
            mirror_mode=False,
            output_file=None,
            project_path=sandbox_dir
        )
        app.push_screen(pscreen)
        await pilot.pause()
        
        all_failed_results = [
            ParseResult(sub_file, ".py", "", "Failed", "Critical Error")
        ]
        pscreen.finished_parsing(all_failed_results)
        await pilot.pause()
        await pilot.click("#close_btn")
        
        # 10. Test Directory Tree Expand & calculate_metrics_async
        tree = app.query_one("#dir_tree")
        for node in tree.root.children:
            if node.data and node.data.path.is_dir():
                tree.post_message(tree.NodeExpanded(node))
                await pilot.pause()
                break

        # Quit
        await pilot.press("q")

