import pytest
from pathlib import Path

@pytest.fixture
def sandbox_dir(tmp_path):
    """Creates a temporary project directory populated with various valid and excluded test files."""
    # 1. Create a Python code file
    code_path = tmp_path / "calculator.py"
    with open(code_path, "w", encoding="utf-8") as f:
        f.write('''def add(a, b):
    return a + b
''')

    # 2. Create a Markdown file
    md_path = tmp_path / "readme.md"
    with open(md_path, "w", encoding="utf-8") as f:
        f.write('''# Readme
- Item A
- Item B
''')

    # 3. Create a DOCX file using python-docx
    import docx
    doc = docx.Document()
    doc.add_heading("Heading 1 Test", level=1)
    doc.add_paragraph("This is paragraph 1 of the Word document.")
    doc.save(tmp_path / "sample.docx")

    # 4. Create a PPTX file using python-pptx
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "RTX TUI Presentation"
    slide.placeholders[1].text = "Bullet Point 1\nBullet Point 2"
    prs.save(tmp_path / "sample.pptx")

    # 5. Create a PDF file using PyMuPDF (fitz)
    import fitz
    from PIL import Image
    import io
    doc_pdf = fitz.open()
    page = doc_pdf.new_page()
    page.insert_text((50, 50), "Hello PDF World from RTX Parser!", fontsize=14)
    
    # Create a simple red image
    img = Image.new("RGB", (100, 100), color="red")
    img_bytes = io.BytesIO()
    img.save(img_bytes, format="PNG")
    
    # Insert image on the page
    page.insert_image(fitz.Rect(100, 100, 200, 200), stream=img_bytes.getvalue())
    doc_pdf.save(tmp_path / "sample.pdf")

    # 6. Create an EPUB file using ebooklib
    import ebooklib
    from ebooklib import epub
    book = epub.EpubBook()
    book.set_identifier("id_rtx_test_123")
    book.set_title("RTX Ebook Test")
    book.set_language("en")
    book.add_author("Developer Antigravity")
    
    c1 = epub.EpubHtml(title="Introduction", file_name="intro.xhtml", lang="en")
    c1.content = "<h1>Introduction to RTX</h1><p>Epub is a great digital book format.</p>"
    book.add_item(c1)
    book.toc = (epub.Link("intro.xhtml", "Introduction", "intro"),)
    book.spine = ["nav", c1]
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    epub.write_epub(tmp_path / "sample.epub", book)

    # 7. Create a CSV file
    csv_path = tmp_path / "sample.csv"
    with open(csv_path, "w", encoding="utf-8", newline="") as f:
        import csv
        writer = csv.writer(f)
        writer.writerow(["Header1", "Header2"])
        writer.writerow(["Row1Col1", "Row1Col2"])
        writer.writerow(["Row2Col1", "Row2Col2"])

    # 8. Create an XLSX file using openpyxl
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "TestSheet"
    ws.append(["ColA", "ColB"])
    ws.append(["ValA1", "ValB1"])
    ws.append(["ValA2", "ValB2"])
    wb.save(tmp_path / "sample.xlsx")

    # 9. Excluded folders
    excluded_dir = tmp_path / "node_modules"
    excluded_dir.mkdir(exist_ok=True)
    with open(excluded_dir / "index.js", "w") as f:
        f.write("console.log('should be ignored');")

    hidden_dir = tmp_path / ".git"
    hidden_dir.mkdir(exist_ok=True)
    with open(hidden_dir / "config", "w") as f:
        f.write("should be ignored")

    binary_file = tmp_path / "image.png"
    with open(binary_file, "wb") as f:
        f.write(b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00")

    null_byte_file = tmp_path / "unknown.dat"
    with open(null_byte_file, "wb") as f:
        f.write(b"Normal text with null byte \x00 in it.")

    return tmp_path
