from pathlib import Path
from rtx.parsers.base import BaseParser

class PptxParser(BaseParser):
    def parse(self, path: Path) -> str:
        # Lazy import to avoid loading heavy modules on startup
        from pptx import Presentation
        
        prs = Presentation(path)
        markdown_lines = []
        
        for idx, slide in enumerate(prs.slides, 1):
            markdown_lines.append(f"\n## Slide {idx}\n")
            
            # Find title if available
            title_text = ""
            if slide.shapes.title:
                title_text = slide.shapes.title.text.strip()
                if title_text:
                    markdown_lines.append(f"### {title_text}\n")
            
            # Extract text from shapes
            for shape in slide.shapes:
                # Skip title as we already processed it
                if shape == slide.shapes.title:
                    continue
                    
                if not shape.has_text_frame:
                    continue
                    
                # Process paragraphs inside text frame
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if not text:
                        continue
                    
                    # Indent depending on level
                    indent = "  " * paragraph.level
                    if paragraph.level > 0 or len(shape.text_frame.paragraphs) > 1:
                        markdown_lines.append(f"{indent}* {text}")
                    else:
                        markdown_lines.append(text)
            
        return "\n".join(markdown_lines).strip()
