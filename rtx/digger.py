import os
import gc
from pathlib import Path
from typing import Dict, Tuple, Set, Optional, Iterator

# Excluded directory names
EXCLUDED_DIR_NAMES = {
    "node_modules",
    ".venv",
    "venv",
    "env",
    "__pycache__",
    ".git",
    ".idea",
    ".vscode",
    ".cache",
    ".rtx",
    "build",
    "dist",
}

# Binary extensions to skip (that aren't our explicitly supported document formats)
BINARY_EXTENSIONS = {
    # Images
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp", ".ico", ".svg",
    # Audio
    ".mp3", ".wav", ".ogg", ".flac", ".m4a", ".aac",
    # Video
    ".mp4", ".mkv", ".avi", ".mov", ".flv", ".wmv", ".webm",
    # Archives
    ".zip", ".tar", ".gz", ".bz2", ".xz", ".7z", ".rar", ".tgz",
    # Binaries/Libraries
    ".exe", ".dll", ".so", ".dylib", ".bin", ".class", ".pyc", ".pyo", ".pyd", ".o", ".a", ".lib", ".out", ".app",
    # Fonts
    ".woff", ".woff2", ".ttf", ".otf", ".eot",
    # Databases
    ".db", ".sqlite", ".sqlite3", ".mdb",
}

# Document extensions we explicitly support
DOCUMENT_EXTENSIONS = {
    ".pdf", ".docx", ".pptx", ".epub", ".fb2"
}

# Programming language code extensions we support (10+ languages)
CODE_EXTENSIONS = {
    ".py", ".go", ".rs", ".c", ".cpp", ".cc", ".h", ".hpp", ".js", ".jsx", ".ts", ".tsx",
    ".html", ".htm", ".css", ".yaml", ".yml", ".json", ".toml", ".sh", ".bash", ".md"
}

def is_binary_content(path: Path) -> bool:
    """Check if a file contains binary content by reading the first 1024 bytes."""
    try:
        with open(path, "rb") as f:
            chunk = f.read(1024)
            return b"\x00" in chunk
    except OSError:
        return True

def is_excluded_path(path: Path, root_path: Path) -> bool:
    """
    Determine if a path should be excluded according to the RTX specifications.
    Checks directory names and binary extensions/contents.
    """
    try:
        # Check if the path itself or any of its parents (relative to root) are in excluded dirs
        try:
            rel_path = path.relative_to(root_path)
        except ValueError:
            return True  # Exclude paths outside the project root

        for part in rel_path.parts:
            # Skip hidden folders/files starting with '.'
            if part.startswith(".") and part not in (".", ".."):
                # Exception: unless it's a file inside a path we explicitly allow, 
                # but standard rules say exclude hidden folders/files.
                return True
            if part in EXCLUDED_DIR_NAMES:
                return True

        if path.is_file():
            suffix = path.suffix.lower()
            if suffix in BINARY_EXTENSIONS:
                return True
            # If it's not a known document or code suffix, do a binary content check
            if suffix not in DOCUMENT_EXTENSIONS and suffix not in CODE_EXTENSIONS:
                if is_binary_content(path):
                    return True

        return False
    except Exception:
        # If any exception occurs during check, play safe and exclude it
        return True

def count_pdf_metrics(path: Path) -> Tuple[int, int]:
    """Calculate lines and characters for a PDF file using pymupdf or pypdf (fast)."""
    text = ""
    # Try pymupdf first
    try:
        import fitz
        doc = fitz.open(path)
        for page in doc:
            t = page.get_text()
            if t:
                text += t
        lines = text.count("\n") + 1 if text else 0
        return lines, len(text)
    except Exception:
        pass

    # Try pypdf fallback
    try:
        from pypdf import PdfReader
        reader = PdfReader(path)
        for page in reader.pages:
            t = page.extract_text()
            if t:
                text += t
        lines = text.count("\n") + 1 if text else 0
        return lines, len(text)
    except Exception:
        return 0, 0

def count_docx_metrics(path: Path) -> Tuple[int, int]:
    """Calculate lines and characters for a DOCX file."""
    try:
        import docx
        doc = docx.Document(path)
        text_parts = [p.text for p in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text_parts.append(cell.text)
        text = "\n".join(text_parts)
        lines = text.count("\n") + 1 if text else 0
        return lines, len(text)
    except Exception:
        return 0, 0

def count_pptx_metrics(path: Path) -> Tuple[int, int]:
    """Calculate lines and characters for a PPTX file."""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text)
        text = "\n".join(text_parts)
        lines = text.count("\n") + 1 if text else 0
        return lines, len(text)
    except Exception:
        return 0, 0

def count_epub_metrics(path: Path) -> Tuple[int, int]:
    """Calculate lines and characters for an EPUB file."""
    try:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup
        
        # Suppress ebooklib warnings
        import warnings
        warnings.filterwarnings("ignore", category=UserWarning)
        
        book = epub.read_epub(path)
        text_parts = []
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), "html.parser")
                text_parts.append(soup.get_text())
        text = "\n".join(text_parts)
        lines = text.count("\n") + 1 if text else 0
        return lines, len(text)
    except Exception:
        return 0, 0

def count_fb2_metrics(path: Path) -> Tuple[int, int]:
    """Calculate lines and characters for an FB2 file."""
    try:
        from bs4 import BeautifulSoup
        with open(path, "rb") as f:
            soup = BeautifulSoup(f.read(), "xml")
            text = soup.get_text()
        lines = text.count("\n") + 1 if text else 0
        return lines, len(text)
    except Exception:
        return 0, 0

def get_file_metrics(path: Path) -> Tuple[int, int, int]:
    """
    Calculate metrics for a single file: (Lines, Chars, Bytes).
    Handles documents and code files dynamically.
    """
    if not path.is_file():
        return 0, 0, 0
    
    bytes_size = path.stat().st_size
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        lines, chars = count_pdf_metrics(path)
    elif suffix == ".docx":
        lines, chars = count_docx_metrics(path)
    elif suffix == ".pptx":
        lines, chars = count_pptx_metrics(path)
    elif suffix == ".epub":
        lines, chars = count_epub_metrics(path)
    elif suffix == ".fb2":
        lines, chars = count_fb2_metrics(path)
    else:
        # Regular text or code files
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
                lines = content.count("\n") + 1 if content else 0
                chars = len(content)
        except Exception:
            lines, chars = 0, 0
            
    return lines, chars, bytes_size

class DiggerEngine:
    def __init__(self, root_path: Path):
        self.root_path = root_path.resolve()
        self.metrics_cache: Dict[Path, Tuple[int, int, int, int]] = {}  # Path -> (Files, Lines, Chars, Bytes)

    def scan_valid_files(self, start_dir: Optional[Path] = None) -> Iterator[Path]:
        """Generator yielding all valid (non-excluded) files under the directory."""
        target_dir = start_dir.resolve() if start_dir else self.root_path
        
        for dirpath, dirnames, filenames in os.walk(target_dir):
            path_dir = Path(dirpath)
            
            # Filter directories in place to prevent os.walk from entering excluded dirs
            # We must iterate over a copy of dirnames since we modify it in-place
            for dirname in list(dirnames):
                sub_dir = path_dir / dirname
                if is_excluded_path(sub_dir, self.root_path):
                    dirnames.remove(dirname)

            for filename in filenames:
                file_path = path_dir / filename
                if not is_excluded_path(file_path, self.root_path):
                    yield file_path

    def calculate_metrics(self, path: Path) -> Tuple[int, int, int, int]:
        """
        Calculate metrics for a path (file or directory).
        Returns a tuple: (Files, Lines, Chars, Bytes).
        For files: (1, Lines, Chars, Bytes).
        Caches results for efficiency.
        """
        path = path.resolve()
        if path in self.metrics_cache:
            return self.metrics_cache[path]

        if path.is_file():
            if is_excluded_path(path, self.root_path):
                metrics = (0, 0, 0, 0)
            else:
                lines, chars, bytes_size = get_file_metrics(path)
                metrics = (1, lines, chars, bytes_size)
            self.metrics_cache[path] = metrics
            return metrics

        # Directory calculation
        total_files = 0
        total_lines = 0
        total_chars = 0
        total_bytes = 0

        # We only walk immediate children or recurse. We use os.scandir for fast single level stats.
        try:
            for entry in os.scandir(path):
                entry_path = Path(entry.path)
                if is_excluded_path(entry_path, self.root_path):
                    continue

                if entry.is_file():
                    lines, chars, bytes_size = get_file_metrics(entry_path)
                    total_files += 1
                    total_lines += lines
                    total_chars += chars
                    total_bytes += bytes_size
                elif entry.is_dir():
                    # Recurse
                    sub_files, sub_lines, sub_chars, sub_bytes = self.calculate_metrics(entry_path)
                    total_files += sub_files
                    total_lines += sub_lines
                    total_chars += sub_chars
                    total_bytes += sub_bytes
        except OSError:
            pass

        metrics = (total_files, total_lines, total_chars, total_bytes)
        self.metrics_cache[path] = metrics
        return metrics

def format_bytes(bytes_count: int) -> str:
    """Helper to convert bytes into a human-readable string (KiB, MiB)."""
    if bytes_count < 1024:
        return f"{bytes_count} B"
    elif bytes_count < 1024 * 1024:
        return f"{bytes_count / 1024:.1f} KiB"
    else:
        return f"{bytes_count / (1024 * 1024):.1f} MiB"
